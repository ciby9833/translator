# backend/services/document_processor.py 
# 文档处理识别不同类型文件器
from typing import List, Optional
import io
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
import markdown
import bs4
import os
import tempfile

class DocumentProcessor:
    """文档处理器：支持不同格式文档的文本提取"""
    
    def __init__(self):
        # 支持的文件扩展名映射
        self.supported_extensions = {
            '.txt': self.process_text,
            '.md': self.process_markdown,
            '.docx': self.process_docx,
            '.pptx': self.process_pptx,
            '.pdf': self.process_pdf
        }

    def can_process(self, filename: str) -> bool:
        """检查是否支持处理该文件类型"""
        ext = os.path.splitext(filename.lower())[1]
        return ext in self.supported_extensions

    def process_file(self, file_path: str) -> str:
        """处理文件并返回提取的文本"""
        ext = os.path.splitext(file_path.lower())[1]
        if ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {ext}")
            
        return self.supported_extensions[ext](file_path)

    def process_text(self, file_path: str) -> str:
        """处理纯文本文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()

    def process_markdown(self, file_path: str) -> str:
        """处理 Markdown 文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            md_text = f.read()
            # 转换 Markdown 为 HTML
            html = markdown.markdown(md_text)
            # 使用 BeautifulSoup 提取纯文本
            soup = bs4.BeautifulSoup(html, 'html.parser')
            return soup.get_text(separator='\n\n')

    def process_docx(self, file_path: str) -> str:
        """处理 Word 文档"""
        doc = Document(file_path)
        text_parts = []
        
        # 处理段落
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # 处理表格
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        return '\n\n'.join(text_parts)

    def process_pptx(self, file_path: str) -> str:
        """处理 PowerPoint 文档"""
        prs = Presentation(file_path)
        text_parts = []
        
        for slide in prs.slides:
            slide_texts = []
            
            # 处理形状中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
                    
                # 处理表格
                if shape.has_table:
                    table_texts = []
                    for row in shape.table.rows:
                        row_text = ' | '.join(
                            cell.text.strip() 
                            for cell in row.cells 
                            if cell.text.strip()
                        )
                        if row_text:
                            table_texts.append(row_text)
                    if table_texts:
                        slide_texts.append('\n'.join(table_texts))
            
            if slide_texts:
                text_parts.append('\n'.join(slide_texts))
        
        return '\n\n'.join(text_parts)

    def process_pdf(self, file_path: str) -> str:
        """处理 PDF 文档"""
        reader = PdfReader(file_path)
        text_parts = []
        
        for page in reader.pages:
            text = page.extract_text()
            if text.strip():
                text_parts.append(text)
        
        return '\n\n'.join(text_parts)

    def get_document_metadata(self, file_path: str) -> dict:
        """获取文档元数据"""
        ext = os.path.splitext(file_path.lower())[1]
        metadata = {
            'file_size': os.path.getsize(file_path),
            'file_type': ext[1:],  # 移除点号
        }
        
        try:
            if ext == '.pdf':
                reader = PdfReader(file_path)
                metadata.update({
                    'pages': len(reader.pages),
                    'info': reader.metadata
                })
            elif ext == '.docx':
                doc = Document(file_path)
                metadata.update({
                    'paragraphs': len(doc.paragraphs),
                    'sections': len(doc.sections)
                })
            elif ext == '.pptx':
                prs = Presentation(file_path)
                metadata.update({
                    'slides': len(prs.slides)
                })
        except Exception as e:
            metadata['error'] = str(e)
            
        return metadata

    async def process_file_async(self, content: bytes, filename: str) -> str:
        """异步处理文件内容"""
        # 根据文件扩展名选择处理方法
        ext = os.path.splitext(filename.lower())[1]
        
        if ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {ext}")
            
        # 创建临时文件
        with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as temp_file:
            temp_file.write(content)
            temp_path = temp_file.name
        
        try:
            # 使用现有的处理方法
            text = self.supported_extensions[ext](temp_path)
            return text
        finally:
            # 清理临时文件
            os.unlink(temp_path)